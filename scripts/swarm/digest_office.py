"""Extract text, notes, and images from Office source material.

Runs at zero LLM cost. Everything downstream reads this output rather
than the original binaries.
"""

from __future__ import annotations

import json
from dataclasses import dataclass, field
from pathlib import Path

import docx
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from swarm.paths import validate_session_id


@dataclass
class Slide:
    index: int
    title: str
    body: str
    notes: str


@dataclass
class DigestResult:
    sid: str
    slides: list[Slide] = field(default_factory=list)
    images: list[dict] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)


def extract_pptx(src: Path, sid: str, out_dir: Path) -> DigestResult:
    """Pull every slide's text, notes, and embedded images."""
    validate_session_id(sid)
    out_dir.mkdir(parents=True, exist_ok=True)
    result = DigestResult(sid=sid)
    prs = Presentation(str(src))

    for i, slide in enumerate(prs.slides, start=1):
        title = ""
        body_parts: list[str] = []

        has_image = False
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or getattr(shape, "image", None) is not None:
                has_image = True
                _save_image(shape, i, len(result.images) + 1, out_dir, result)
                continue
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if shape == slide.shapes.title:
                title = text
            else:
                body_parts.append(text)

        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        if not title and not body_parts and not has_image:
            result.warnings.append(f"slide {i} is empty (no title, no body)")

        result.slides.append(
            Slide(index=i, title=title, body="\n\n".join(body_parts), notes=notes)
        )

    _write_manifest(out_dir, result)
    return result


def _save_image(shape, slide_index: int, seq: int, out_dir: Path, result: DigestResult) -> None:
    try:
        image = shape.image
    except (AttributeError, ValueError) as exc:
        result.warnings.append(f"slide {slide_index}: failed to save image ({exc})")
        return
    name = f"img-{seq:02d}.{image.ext}"
    (out_dir / name).write_bytes(image.blob)
    result.images.append(
        {
            "file": name,
            "slide": slide_index,
            "ext": image.ext,
            "bytes": len(image.blob),
        }
    )


def _write_manifest(out_dir: Path, result: DigestResult) -> None:
    (out_dir / "manifest.json").write_text(
        json.dumps({"id": result.sid, "images": result.images}, indent=2),
        encoding="utf-8",
    )


def extract_docx(src: Path, sid: str) -> str:
    """Return a docx's paragraphs as markdown-ish plain text."""
    validate_session_id(sid)
    document = docx.Document(str(src))
    return "\n\n".join(p.text.strip() for p in document.paragraphs if p.text.strip())
