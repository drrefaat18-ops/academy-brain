"""Composite reserved regions into an exported deck, then prove none was missed.

§5b of `docs/PLAN-academy-template-and-ev3.md`. The L1-s1 run reserved regions for
assets it already had on disk and then never composited them in. From the owner's
seat an empty dashed box is indistinguishable from a demand for content he cannot
supply — he has no way to know the file is already there. That is precisely what
the owner rule exists to prevent, and it reached him anyway.

Reserving a region is half the job. This is the other half.

A `[Reserved Image Area: …]` box visible in a delivered artifact is an unfinished
build, not a deliverable. `assert_filled` fails closed so it can never be reached.

Two formats, because the pipeline has two. `generate_session.py` downloads a **PDF**
(`_run_pass`: `out = out_dir / f"{ps.key}.pdf"`), so the PDF path is the one that runs
in production. PPTX is supported because the deck is edited as a deck before delivery
and the same contract has to hold there. Format is chosen by suffix; anything else is
refused rather than guessed at.
"""

from __future__ import annotations

import re
from dataclasses import dataclass
from pathlib import Path

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Emu

# What the generator writes into the prompt, and therefore what NotebookLM renders.
REGION = re.compile(r"\[Reserved Image Area:\s*([^\]]+?)\s*\]")


class OverlayError(RuntimeError):
    """The overlay could not be completed. Never warned about, never worked around."""


@dataclass(frozen=True)
class Region:
    slide_index: int
    aid: str
    left: int
    top: int
    width: int
    height: int


def _find_regions_pptx(deck: Path) -> list[Region]:
    prs = Presentation(str(deck))
    found: list[Region] = []
    for i, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for aid in REGION.findall(shape.text_frame.text):
                found.append(
                    Region(
                        slide_index=i,
                        aid=aid.strip(),
                        left=shape.left,
                        top=shape.top,
                        width=shape.width,
                        height=shape.height,
                    )
                )
    return found


def _fit(box_w: int, box_h: int, img_w: int, img_h: int) -> tuple[int, int]:
    """Largest size fitting the box without distorting the image.

    A stretched screenshot of code is a wrong screenshot — the same failure the
    EVIDENCE class exists to prevent, arriving through the back door.
    """
    if img_w <= 0 or img_h <= 0:
        raise OverlayError("image reports a non-positive dimension")
    scale = min(box_w / img_w, box_h / img_h)
    return int(img_w * scale), int(img_h * scale)


def _remove(shape) -> None:
    shape._element.getparent().remove(shape._element)


def _overlay_pptx(deck: Path, assets: dict[str, Path], out: Path | None) -> list[str]:
    prs = Presentation(str(deck))
    filled: list[str] = []
    missing: list[str] = []

    for i, slide in enumerate(prs.slides, start=1):
        for shape in list(slide.shapes):
            if not shape.has_text_frame:
                continue
            aids = REGION.findall(shape.text_frame.text)
            if not aids:
                continue
            if len(aids) > 1:
                raise OverlayError(
                    f"slide {i}: one shape declares {len(aids)} reserved regions "
                    f"({', '.join(a.strip() for a in aids)}) — regions must not collide"
                )
            aid = aids[0].strip()
            img = assets.get(aid)
            if img is None or not Path(img).is_file():
                missing.append(f"slide {i}: {aid}")
                continue

            box = (shape.left, shape.top, shape.width, shape.height)
            pic = slide.shapes.add_picture(str(img), box[0], box[1])
            w, h = _fit(box[2], box[3], pic.width, pic.height)
            pic.width, pic.height = w, h
            # centre it in the reserved box rather than pinning it to a corner
            pic.left = box[0] + (box[2] - w) // 2
            pic.top = box[1] + (box[3] - h) // 2
            _remove(shape)
            filled.append(aid)

    if missing:
        raise OverlayError(
            "reserved regions with no resolvable asset: "
            + "; ".join(missing)
            + ". Resolve this BEFORE generating — redesign, substitute, or drop. "
            "Do not deliver a deck with an empty box in it."
        )

    prs.save(str(out or deck))
    return filled


# --------------------------------------------------------------------------
# PDF — the format production actually downloads
# --------------------------------------------------------------------------


def _find_regions_pdf(deck: Path) -> list[Region]:
    found: list[Region] = []
    with fitz.open(str(deck)) as doc:
        for i, page in enumerate(doc, start=1):
            for aid in REGION.findall(page.get_text()):
                aid = aid.strip()
                rects = page.search_for(f"[Reserved Image Area: {aid}]")
                if not rects:
                    # the marker is in the text layer but not locatable as a run —
                    # record it with a zero box so assert_filled still catches it
                    found.append(Region(i, aid, 0, 0, 0, 0))
                    continue
                for r in rects:
                    found.append(
                        Region(i, aid, int(r.x0), int(r.y0), int(r.width), int(r.height))
                    )
    return found


# Which asset ids this module actually placed, written into the PDF so the proof
# survives the process that made it. Counting images cannot do this job: a deck
# carrying decorative chrome already has images, so an interrupted run that
# redacted every marker and inserted nothing still cleared a count-based check.
_PLACED_PREFIX = "swarm-overlay-placed:"


def _record_placements(doc, placements: dict[str, tuple[int, int]]) -> None:
    """Record aid=page:xref for each placement, not merely the aid.

    An id alone says an insert once happened. It does not say the image is still
    there, and a record that outlives what it describes is a claim rather than a
    proof — the whole defect class this module keeps rediscovering.
    """
    meta = dict(doc.metadata or {})
    keywords = [
        k for k in (meta.get("keywords") or "").split() if not k.startswith(_PLACED_PREFIX)
    ]
    known = dict(_placements_from(meta))
    known.update(placements)
    body = ",".join(f"{aid}={pg}:{xref}" for aid, (pg, xref) in sorted(known.items()))
    keywords.append(_PLACED_PREFIX + body)
    meta["keywords"] = " ".join(keywords)
    doc.set_metadata(meta)


def _placements_from(meta: dict) -> dict[str, tuple[int, int]]:
    for token in (meta.get("keywords") or "").split():
        if not token.startswith(_PLACED_PREFIX):
            continue
        out: dict[str, tuple[int, int]] = {}
        for item in token[len(_PLACED_PREFIX):].split(","):
            aid, _, loc = item.partition("=")
            page, _, xref = loc.partition(":")
            if aid and page.isdigit() and xref.isdigit():
                out[aid] = (int(page), int(xref))
        return out
    return {}


def placements(deck: Path) -> dict[str, tuple[int, int]]:
    """asset id -> (1-based page, image xref) as recorded by :func:`overlay`."""
    doc = fitz.open(str(deck))
    try:
        return _placements_from(dict(doc.metadata or {}))
    finally:
        doc.close()


def placed_ids(deck: Path) -> set[str]:
    """Asset ids with a placement record. Presence of a record is NOT proof."""
    return set(placements(deck))


def _overlay_pdf(deck: Path, assets: dict[str, Path], out: Path | None) -> list[str]:
    filled: list[str] = []
    missing: list[str] = []

    doc = fitz.open(str(deck))
    pending: list[tuple] = []
    placements: dict[str, tuple[int, int]] = {}
    try:
        for i, page in enumerate(doc, start=1):
            for aid in {a.strip() for a in REGION.findall(page.get_text())}:
                rects = page.search_for(f"[Reserved Image Area: {aid}]")
                if not rects:
                    raise OverlayError(
                        f"page {i}: region {aid!r} is in the text but its box cannot be "
                        "located — refusing to guess where the asset goes"
                    )
                if len(rects) > 1:
                    raise OverlayError(
                        f"page {i}: region {aid!r} appears {len(rects)} times — "
                        "regions must not collide"
                    )
                img = assets.get(aid)
                if img is None or not Path(img).is_file():
                    missing.append(f"page {i}: {aid}")
                    continue

                box = rects[0]
                # REDACT, do not paint over. A white rectangle hides the marker
                # visually but leaves it in the text layer, so assert_filled would
                # still report the region unfilled on a deck that looks correct —
                # and a text search of the delivered PDF would still find it.
                page.add_redact_annot(box)
                pending.append((page, box, img, aid))

        # redactions apply per page, and applying them invalidates rects, so do it
        # once per page before any image is placed
        for page in {id(pg): pg for pg, *_ in pending}.values():
            page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_NONE)

        for page, box, img, aid in pending:
            # keep_proportion is PyMuPDF's aspect-preserving fit — a stretched
            # screenshot of code is a wrong screenshot.
            before = {x[0] for x in page.get_images(full=True)}
            page.insert_image(box, filename=str(img), keep_proportion=True)
            new_xrefs = {x[0] for x in page.get_images(full=True)} - before
            if not new_xrefs:
                raise OverlayError(
                    f"page {page.number + 1}: inserting {aid} added no image to the "
                    "page. Refusing to record a placement that did not happen."
                )
            placements[aid] = (page.number + 1, min(new_xrefs))
            filled.append(aid)

        if missing:
            raise OverlayError(
                "reserved regions with no resolvable asset: "
                + "; ".join(missing)
                + ". Resolve this BEFORE generating — redesign, substitute, or drop. "
                "Do not deliver a deck with an empty box in it."
            )

        _record_placements(doc, placements)

        target = Path(out) if out else Path(deck)
        if target == Path(deck):
            doc.saveIncr()
        else:
            doc.save(str(target))
    finally:
        doc.close()

    return filled


# --------------------------------------------------------------------------
# front door — format by suffix, never guessed
# --------------------------------------------------------------------------

_HANDLERS = {
    ".pdf": (_find_regions_pdf, _overlay_pdf),
    ".pptx": (_find_regions_pptx, _overlay_pptx),
}


def _handlers(deck: Path):
    suffix = deck.suffix.lower()
    if suffix not in _HANDLERS:
        raise OverlayError(
            f"unsupported deck format {suffix!r} for {deck} — expected one of "
            f"{sorted(_HANDLERS)}"
        )
    return _HANDLERS[suffix]


def find_regions(deck: Path) -> list[Region]:
    """Every reserved region still present in the deck, in page/slide order."""
    deck = Path(deck)
    if not deck.is_file():
        raise OverlayError(f"no deck at {deck}")
    return _handlers(deck)[0](deck)


def overlay(deck: Path, assets: dict[str, Path], out: Path | None = None) -> list[str]:
    """Composite each mapped asset into its reserved region. Returns filled asset ids.

    `assets` maps asset id -> image path. An id present in the deck but absent from
    the mapping is an error, not a skip: silently leaving a region empty is the
    original defect this module exists to close.
    """
    deck = Path(deck)
    if not deck.is_file():
        raise OverlayError(f"no deck at {deck}")
    return _handlers(deck)[1](deck, assets, out)


def images_on_page(deck: Path) -> dict[int, int]:
    """Count of embedded images per 1-based page. PDF only."""
    counts: dict[int, int] = {}
    with fitz.open(str(deck)) as doc:
        for i, page in enumerate(doc, start=1):
            counts[i] = len(page.get_images())
    return counts


def assert_filled(deck: Path, expected: dict[str, Path] | None = None) -> None:
    """Fail closed if the deck is not actually finished.

    Marker ABSENCE is not proof of image PRESENCE, and treating it as such was
    this module's own version of the bug it exists to catch. An interrupted run
    leaves the redaction applied and the image never inserted: no marker, no
    picture, and the old check called that a pass. Verified by construction —
    a redaction-only PDF passed the previous implementation with zero images
    embedded.

    For a recorded placement, listing the xref in the page's resources is not
    proof either: the exact image object must still be on the page it was
    recorded on AND occupy rendered space there (a non-empty
    ``page.get_image_rects``), which fails for an xref left listed by a
    redaction after its placement was stripped.

    Pass ``expected`` (the asset-id -> path mapping the deck was built from) to
    get the real check. Without it only the weaker marker scan runs, which is
    all a PPTX deck supports.
    """
    deck = Path(deck)
    left = find_regions(deck)
    if left:
        unit = "page" if deck.suffix.lower() == ".pdf" else "slide"
        where = "; ".join(f"{unit} {r.slide_index}: {r.aid}" for r in left)
        raise OverlayError(
            f"{len(left)} reserved region(s) still unfilled in {deck}: {where}. "
            "This deck is an unfinished build, not a deliverable. It must not reach "
            "the owner — an empty dashed box reads to him as a demand for content he "
            "cannot supply, when the asset is usually already on disk."
        )

    if not expected or deck.suffix.lower() != ".pdf":
        return

    recorded = placements(deck)
    unrecorded = sorted(set(expected) - set(recorded))

    # A record is a claim about the past. Check the present: the exact image
    # object each placement named must still be on the page it named. Trusting
    # the record alone would repeat this module's own defect one level higher —
    # it was already caught doing that with a marker, and then with a count.
    doc = fitz.open(str(deck))
    try:
        gone = []
        for aid in sorted(set(expected) & set(recorded)):
            page_no, xref = recorded[aid]
            if page_no < 1 or page_no > doc.page_count:
                gone.append(f"{aid} (recorded on page {page_no}, which does not exist)")
                continue
            page = doc[page_no - 1]
            live = {x[0] for x in page.get_images(full=True)}
            if xref not in live:
                gone.append(f"{aid} (image {xref} no longer on page {page_no})")
                continue
            # Listed is not placed. A redacting writer can strip every draw
            # command for an image and still leave the xref in the page's
            # resources — membership alone was checked once before and passed a
            # deck whose asset occupied no rendered space at all.
            if not page.get_image_rects(xref):
                gone.append(
                    f"{aid} (image {xref} listed on page {page_no} but has no "
                    "placement rect)"
                )
    finally:
        doc.close()

    if unrecorded or gone:
        total = sum(images_on_page(deck).values())
        problems = []
        if unrecorded:
            problems.append(f"never recorded as placed: {', '.join(unrecorded)}")
        if gone:
            problems.append("recorded but no longer present: " + "; ".join(gone))
        raise OverlayError(
            f"{deck} has no marker left, but " + "; ".join(problems) + f". The deck "
            f"carries {total} embedded image(s) in total, which proves nothing — "
            "chrome and mascots are images too. A deck with the placeholder removed "
            "and the picture missing is the worst outcome available: it looks "
            "finished and is not. Delete it and regenerate."
        )


def _demo() -> None:
    import tempfile

    from pptx import Presentation as P

    with tempfile.TemporaryDirectory() as td:
        d = Path(td)

        # a 1x1 PNG, enough for python-pptx to embed
        png = d / "bug1.png"
        png.write_bytes(
            bytes.fromhex(
                "89504e470d0a1a0a0000000d49484452000000010000000108060000001f15c4"
                "890000000a49444154789c6360000002000100ffff03000006000557bfabd400"
                "00000049454e44ae426082"
            )
        )

        deck = d / "deck.pptx"
        prs = P()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        box = slide.shapes.add_textbox(Emu(914400), Emu(914400), Emu(2743200), Emu(1828800))
        box.text_frame.text = "[Reserved Image Area: bug-1]"
        prs.save(str(deck))

        regions = find_regions(deck)
        assert len(regions) == 1 and regions[0].aid == "bug-1", regions

        # an unfilled deck must fail closed
        try:
            assert_filled(deck)
        except OverlayError as exc:
            assert "bug-1" in str(exc)
        else:  # pragma: no cover
            raise AssertionError("assert_filled passed a deck with an empty region")

        # a region whose asset is not mapped is an error, never a silent skip
        try:
            overlay(deck, {}, out=d / "x.pptx")
        except OverlayError as exc:
            assert "no resolvable asset" in str(exc), exc
        else:  # pragma: no cover
            raise AssertionError("overlay silently skipped an unmapped region")

        out = d / "filled.pptx"
        assert overlay(deck, {"bug-1": png}, out=out) == ["bug-1"]
        assert find_regions(out) == []
        assert_filled(out)

        # the placeholder is gone and a picture stands where it was
        pics = [s for s in P(str(out)).slides[0].shapes if s.shape_type == 13]
        assert len(pics) == 1, pics

    print("overlay.py self-check OK")


if __name__ == "__main__":
    _demo()
