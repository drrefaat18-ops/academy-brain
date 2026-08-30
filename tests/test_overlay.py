"""§5b overlay: the half of the reserved-region contract that was never built.

The regression these guard: L1-s1 reserved regions for assets that were already on
disk, then shipped the deck with empty dashed boxes in it.
"""

import pytest
from pptx import Presentation
from pptx.util import Emu

from swarm import overlay

def _png(tmp_path, name="img.png", size=(64, 64)):
    """A real, decodable PNG.

    A hand-rolled 1x1 byte string is not enough: python-pptx embeds image bytes
    without decoding them, but MuPDF actually parses the image, so a malformed
    fixture passes the PPTX tests and fails only on the PDF path — which is the
    path production uses.
    """
    from PIL import Image

    p = tmp_path / name
    Image.new("RGB", size, (200, 30, 30)).save(p)
    return p


def _deck(tmp_path, *texts, name="deck.pptx"):
    """One slide per text, each carrying that text in a textbox."""
    prs = Presentation()
    for t in texts:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        box = slide.shapes.add_textbox(Emu(914400), Emu(914400), Emu(2743200), Emu(1828800))
        box.text_frame.text = t
    path = tmp_path / name
    prs.save(str(path))
    return path


def test_finds_every_reserved_region_in_slide_order(tmp_path):
    deck = _deck(
        tmp_path,
        "[Reserved Image Area: bug-1]",
        "nothing reserved here",
        "[Reserved Image Area: bug-2]",
    )
    regions = overlay.find_regions(deck)
    assert [(r.slide_index, r.aid) for r in regions] == [(1, "bug-1"), (3, "bug-2")]


def test_unfilled_deck_fails_closed(tmp_path):
    deck = _deck(tmp_path, "[Reserved Image Area: bug-1]")
    with pytest.raises(overlay.OverlayError, match="bug-1"):
        overlay.assert_filled(deck)


def test_clean_deck_passes(tmp_path):
    overlay.assert_filled(_deck(tmp_path, "an ordinary slide"))


def test_overlay_fills_region_and_removes_placeholder(tmp_path):
    deck = _deck(tmp_path, "[Reserved Image Area: bug-1]")
    out = tmp_path / "out.pptx"

    assert overlay.overlay(deck, {"bug-1": _png(tmp_path)}, out=out) == ["bug-1"]

    assert overlay.find_regions(out) == []
    overlay.assert_filled(out)
    pics = [s for s in Presentation(str(out)).slides[0].shapes if s.shape_type == 13]
    assert len(pics) == 1


def test_unmapped_region_is_an_error_not_a_skip(tmp_path):
    """The original defect: an empty region that nothing complained about."""
    deck = _deck(tmp_path, "[Reserved Image Area: bug-1]")
    with pytest.raises(overlay.OverlayError, match="no resolvable asset"):
        overlay.overlay(deck, {}, out=tmp_path / "out.pptx")


def test_mapped_asset_missing_on_disk_is_an_error(tmp_path):
    deck = _deck(tmp_path, "[Reserved Image Area: bug-1]")
    with pytest.raises(overlay.OverlayError, match="no resolvable asset"):
        overlay.overlay(deck, {"bug-1": tmp_path / "gone.png"}, out=tmp_path / "out.pptx")


def test_partial_fill_still_fails_closed(tmp_path):
    """Two regions, one asset: the run must fail, not deliver one empty box."""
    deck = _deck(
        tmp_path, "[Reserved Image Area: bug-1]", "[Reserved Image Area: bug-2]"
    )
    with pytest.raises(overlay.OverlayError, match="bug-2"):
        overlay.overlay(deck, {"bug-1": _png(tmp_path)}, out=tmp_path / "out.pptx")


def test_colliding_regions_in_one_shape_are_refused(tmp_path):
    deck = _deck(tmp_path, "[Reserved Image Area: bug-1] [Reserved Image Area: bug-2]")
    with pytest.raises(overlay.OverlayError, match="must not collide"):
        overlay.overlay(deck, {}, out=tmp_path / "out.pptx")


def test_image_is_not_distorted_to_fill_the_box(tmp_path):
    """A stretched screenshot of code is a wrong screenshot."""
    deck = _deck(tmp_path, "[Reserved Image Area: bug-1]")
    out = tmp_path / "out.pptx"
    overlay.overlay(deck, {"bug-1": _png(tmp_path)}, out=out)

    pic = [s for s in Presentation(str(out)).slides[0].shapes if s.shape_type == 13][0]
    # source png is square, so the placed picture must be square too
    assert pic.width == pic.height


def test_missing_deck_is_an_error(tmp_path):
    with pytest.raises(overlay.OverlayError, match="no deck at"):
        overlay.overlay(tmp_path / "nope.pptx", {})


# --------------------------------------------------------------------------
# PDF — the format _run_pass actually downloads. These are the ones that matter.
# --------------------------------------------------------------------------

import fitz


def _pdf(tmp_path, *texts, name="deck.pdf"):
    """One page per text."""
    doc = fitz.open()
    for t in texts:
        page = doc.new_page()
        page.insert_text((72, 144), t, fontsize=14)
    path = tmp_path / name
    doc.save(str(path))
    doc.close()
    return path


def test_pdf_finds_reserved_regions(tmp_path):
    deck = _pdf(tmp_path, "[Reserved Image Area: bug-1]", "ordinary page")
    regions = overlay.find_regions(deck)
    assert [(r.slide_index, r.aid) for r in regions] == [(1, "bug-1")]
    assert regions[0].width > 0 and regions[0].height > 0


def test_pdf_unfilled_fails_closed(tmp_path):
    deck = _pdf(tmp_path, "[Reserved Image Area: bug-1]")
    with pytest.raises(overlay.OverlayError, match="page 1: bug-1"):
        overlay.assert_filled(deck)


def test_pdf_overlay_fills_and_marker_is_gone(tmp_path):
    deck = _pdf(tmp_path, "[Reserved Image Area: bug-1]")
    out = tmp_path / "out.pdf"

    assert overlay.overlay(deck, {"bug-1": _png(tmp_path)}, out=out) == ["bug-1"]

    assert overlay.find_regions(out) == []
    overlay.assert_filled(out)
    with fitz.open(str(out)) as doc:
        assert len(doc[0].get_images()) == 1


def test_pdf_unmapped_region_is_an_error(tmp_path):
    deck = _pdf(tmp_path, "[Reserved Image Area: bug-1]")
    with pytest.raises(overlay.OverlayError, match="no resolvable asset"):
        overlay.overlay(deck, {}, out=tmp_path / "out.pdf")


def test_pdf_partial_fill_fails_closed(tmp_path):
    deck = _pdf(tmp_path, "[Reserved Image Area: bug-1]", "[Reserved Image Area: bug-2]")
    with pytest.raises(overlay.OverlayError, match="bug-2"):
        overlay.overlay(deck, {"bug-1": _png(tmp_path)}, out=tmp_path / "out.pdf")


def test_unknown_format_is_refused_not_guessed(tmp_path):
    odd = tmp_path / "deck.docx"
    odd.write_bytes(b"not a deck")
    with pytest.raises(overlay.OverlayError, match="unsupported deck format"):
        overlay.find_regions(odd)


def test_missing_deck_is_an_error_for_find_too(tmp_path):
    with pytest.raises(overlay.OverlayError, match="no deck at"):
        overlay.find_regions(tmp_path / "nope.pdf")


# --------------------------------------------------------------------------
# the wiring: generate_session._composite is what makes the module load-bearing.
# Without these, overlay.py is a tested module nothing calls.
# --------------------------------------------------------------------------

from swarm import generate_session as gs


def _pass(*assets):
    return gs.Pass(
        key="deck-a",
        notebook="nb",
        instructions="",
        evidence=[
            gs.Asset(aid=aid, slide="1", path=path, klass="EVIDENCE",
                     status="produced and mapped")
            for aid, path in assets
        ],
    )


def test_composite_fills_the_downloaded_deck(tmp_path):
    deck = _pdf(tmp_path, "[Reserved Image Area: bug-1]")
    gs._composite(deck, _pass(("bug-1", _png(tmp_path))))
    assert overlay.find_regions(deck) == []


def test_composite_refuses_a_deck_it_cannot_fill(tmp_path):
    """The L1-s1 defect: an empty box reaching the owner."""
    deck = _pdf(tmp_path, "[Reserved Image Area: bug-1]")
    with pytest.raises(overlay.OverlayError, match="no resolvable asset"):
        gs._composite(deck, _pass())


def test_composite_passes_a_deck_with_no_regions(tmp_path):
    gs._composite(_pdf(tmp_path, "an ordinary slide"), _pass())


def test_composite_is_idempotent(tmp_path):
    """The skip path re-runs it on an already-overlaid deck."""
    deck = _pdf(tmp_path, "[Reserved Image Area: bug-1]")
    p = _pass(("bug-1", _png(tmp_path)))
    gs._composite(deck, p)
    gs._composite(deck, p)
    assert overlay.find_regions(deck) == []


# --------------------------------------------------------------------------
# C1 review: the gate proved marker ABSENCE, not image PRESENCE — this module's
# own version of the bug it exists to catch.
# --------------------------------------------------------------------------


def _redaction_only(tmp_path, decoys=0):
    """A deck as an interrupted run leaves it: marker redacted, image never placed.

    ``decoys`` embeds unrelated images first — brand chrome, a mascot — which is
    what a real deck carries and what defeated the count-based check.
    """
    deck = _pdf(tmp_path, "[Reserved Image Area: bug-1]")
    doc = fitz.open(str(deck))
    page = doc[0]
    for n in range(decoys):
        page.insert_image(
            fitz.Rect(10 + 30 * n, 10, 35 + 30 * n, 35),
            filename=str(_png(tmp_path, name=f"decoy-{n}.png")),
        )
    page.add_redact_annot(page.search_for("[Reserved Image Area: bug-1]")[0])
    page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_NONE)
    doc.saveIncr()
    doc.close()
    return deck


def test_redaction_without_image_is_caught(tmp_path):
    """C1-01. No marker AND no picture is the worst outcome: it looks finished."""
    deck = _redaction_only(tmp_path)
    assert overlay.find_regions(deck) == []          # marker really is gone
    assert sum(overlay.images_on_page(deck).values()) == 0

    with pytest.raises(overlay.OverlayError, match="no marker left"):
        overlay.assert_filled(deck, {"bug-1": _png(tmp_path)})


def test_a_properly_overlaid_deck_passes_the_stronger_check(tmp_path):
    deck = _pdf(tmp_path, "[Reserved Image Area: bug-1]")
    out = tmp_path / "out.pdf"
    overlay.overlay(deck, {"bug-1": _png(tmp_path)}, out=out)
    overlay.assert_filled(out, {"bug-1": _png(tmp_path)})


def test_composite_rejects_duplicate_asset_ids(tmp_path):
    """C1-06. A dict comprehension let the last duplicate win silently."""
    png = _png(tmp_path)
    p = gs.Pass(
        key="deck-a", notebook="nb", instructions="",
        evidence=[
            gs.Asset(aid="bug-1", slide="1", path=png, klass="EVIDENCE",
                     status="produced and mapped"),
            gs.Asset(aid=" Bug-1 ", slide="2", path=tmp_path / "other.png",
                     klass="EVIDENCE", status="produced and mapped"),
        ],
    )
    with pytest.raises(gs.HardStop, match="share the id"):
        gs._composite(_pdf(tmp_path, "ordinary"), p)


@pytest.mark.parametrize("status", ["unapproved", "not-approved", "NOT APPROVED"])
def test_approval_status_is_exact_not_substring(tmp_path, status):
    """C1-03. Every one of these CONTAINS 'approved' and used to pass the gate."""
    bp = tmp_path / "blueprint.md"
    bp.write_text(
        f"---\nstatus: {status}\napproval:\n  kind: specialist_council\n---\n",
        encoding="utf-8",
    )
    with pytest.raises(gs.HardStop, match="has not been approved"):
        gs.enforce_blueprint_gate(bp)


def test_approval_kind_must_live_in_the_approval_mapping(tmp_path):
    """C1-04. A loose `kind:` search matched unrelated metadata anywhere."""
    bp = tmp_path / "blueprint.md"
    bp.write_text(
        "---\nstatus: approved\nsomething_else:\n  kind: specialist_council\n---\n",
        encoding="utf-8",
    )
    with pytest.raises(gs.HardStop, match="no approval mapping"):
        gs.enforce_blueprint_gate(bp)


@pytest.mark.parametrize(
    "marker",
    ["GAP - owner must decide", "Gap - owner must decide",
     "GAP – owner must decide", "  gap  -  owner   must  decide"],
)
def test_legacy_gap_is_caught_in_every_spelling(tmp_path, marker):
    """C1-05. One exact ASCII form let the other spellings ship as settled."""
    bp = tmp_path / "blueprint.md"
    bp.write_text(
        f"---\nstatus: approved\napproval:\n  kind: specialist_council\n---\n{marker}\n",
        encoding="utf-8",
    )
    with pytest.raises(gs.HardStop, match="untyped marker"):
        gs.enforce_blueprint_gate(bp)


# ---------------------------------------------------------------------------
# Round 5. The count-based proof was itself a fail-open.
# ---------------------------------------------------------------------------


def test_unrelated_images_do_not_prove_the_asset_was_placed(tmp_path):
    """C1-5-01. A deck carrying decorative chrome already has embedded images, so
    an interrupted run that redacted the marker and inserted nothing cleared the
    old total-count check. Presence must be proven per asset."""
    deck = _redaction_only(tmp_path, decoys=3)

    assert overlay.find_regions(deck) == []
    assert sum(overlay.images_on_page(deck).values()) >= 1, "decoys must be present"

    with pytest.raises(overlay.OverlayError, match="never recorded as placed"):
        overlay.assert_filled(deck, {"bug-1": _png(tmp_path)})


def test_overlay_records_every_id_it_placed(tmp_path):
    deck = _pdf(tmp_path, "[Reserved Image Area: bug-1]")
    png = _png(tmp_path)

    assert overlay.overlay(deck, {"bug-1": png}) == ["bug-1"]
    assert overlay.placed_ids(deck) == {"bug-1"}
    overlay.assert_filled(deck, {"bug-1": png})  # must not raise


def test_asset_gate_rejects_duplicate_ids_across_all_rows(tmp_path):
    """C1-5-02. The duplicate check lived only in _evidence_map, so a collision
    between two REFERENCE rows — or across two passes — cleared preflight and was
    settled by whichever row happened to be parsed last."""
    png = _png(tmp_path)
    rows = [
        gs.Asset(aid="img-05", slide="3", path=png, klass="REFERENCE",
                 status="produced and mapped"),
        gs.Asset(aid="IMG-05 ", slide="9", path=png, klass="REFERENCE",
                 status="produced and mapped"),
    ]

    with pytest.raises(gs.HardStop, match="duplicate asset id"):
        gs.enforce_asset_gate(rows)

    assert gs.enforce_asset_gate(rows[:1]) is None


def test_a_placement_record_does_not_survive_the_image_it_describes(tmp_path):
    """C6, C1 lane. The record is a claim about the past. An id-only record made
    a deck whose image was stripped afterwards pass — the same defect as the
    marker and the count, one level up."""
    deck = _pdf(tmp_path, "[Reserved Image Area: bug-1]")
    png = _png(tmp_path)
    overlay.overlay(deck, {"bug-1": png})

    page_no, xref = overlay.placements(deck)["bug-1"]
    assert (page_no, xref) == (1, xref)

    doc = fitz.open(str(deck))
    page = doc[page_no - 1]
    box = page.search_for("bug-1")  # any remaining trace, image itself has no text
    # Strip the image itself by redacting the whole page area with image removal —
    # delete_image() leaves the xref listed (as a blanked placeholder), which
    # would not exercise this check.
    page.add_redact_annot(page.rect)
    page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_REMOVE)
    doc.saveIncr()
    doc.close()

    assert overlay.placed_ids(deck) == {"bug-1"}, "the record survives, as it would in the wild"
    with pytest.raises(overlay.OverlayError, match="no longer on page"):
        overlay.assert_filled(deck, {"bug-1": png})


def test_a_listed_xref_with_no_placement_rect_is_caught(tmp_path):
    """C6 follow-up. Redaction-family tools can leave an image xref LISTED in
    the page's resources while stripping everything that draws it, so resource
    membership alone proves nothing: only a live placement rect does."""
    deck = _pdf(tmp_path, "[Reserved Image Area: bug-1]")
    png = _png(tmp_path)
    overlay.overlay(deck, {"bug-1": png})

    page_no, xref = overlay.placements(deck)["bug-1"]

    # Produce exactly that state: keep the resource entry, remove every draw
    # command. (A whole-page add_redact_annot + apply_redactions(
    # images=fitz.PDF_REDACT_IMAGE_REMOVE) on this PyMuPDF deletes the resource
    # entry outright — the sibling case the test above already covers.)
    doc = fitz.open(str(deck))
    for c in doc[page_no - 1].get_contents():
        doc.update_stream(c, b" ")
    doc.saveIncr()
    doc.close()

    with fitz.open(str(deck)) as proof:
        assert xref in {x[0] for x in proof[page_no - 1].get_images(full=True)}, (
            "fixture must leave the xref listed"
        )
        assert proof[page_no - 1].get_image_rects(xref) == []

    with pytest.raises(overlay.OverlayError, match="no placement rect"):
        overlay.assert_filled(deck, {"bug-1": png})


def test_evidence_clause_emits_the_marker_overlay_parses():
    """L2-s2 regression. The prompt producer and the overlay consumer have to
    agree on one literal marker syntax. They did not: evidence_clause() only
    asked for "a single clean empty image area" in prose, so NotebookLM never
    printed `[Reserved Image Area: <aid>]` and assert_filled() failed the deck
    for assets it could not find a marker for. This asserts the contract
    directly — the clause must contain a string that overlay.REGION matches,
    carrying the real asset id."""
    from pathlib import Path

    from swarm.generate_session import Asset, evidence_clause

    ev = [
        Asset(
            aid="ev3-l2s2-img02",
            slide="5",
            path=Path("ev3-l2s2-img02.png"),
            klass="EVIDENCE",
            status="Produced and mapped",
        )
    ]
    clause = evidence_clause(ev)

    assert overlay.REGION.findall(clause) == ["ev3-l2s2-img02"], (
        "evidence_clause must emit a marker overlay.REGION can parse"
    )
    assert "[Reserved Image Area: ev3-l2s2-img02]" in clause
