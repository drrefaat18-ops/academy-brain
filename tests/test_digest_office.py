from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from swarm import digest_office


@pytest.fixture
def sample_pptx(tmp_path: Path) -> Path:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "What is a micro:bit?"
    slide.placeholders[1].text = "A tiny computer you can program."
    slide.notes_slide.notes_text_frame.text = "Ask students to predict first."

    second = prs.slides.add_slide(prs.slide_layouts[5])
    second.shapes.title.text = "The LED grid"
    png = tmp_path / "dot.png"
    png.write_bytes(
        bytes.fromhex(
            "89504e470d0a1a0a0000000d4948445200000001000000010802000000907753"
            "de0000000c4944415408d76360000000020001e221bc330000000049454e44ae426082"
        )
    )
    second.shapes.add_picture(str(png), Inches(1), Inches(1))

    out = tmp_path / "session-1.pptx"
    prs.save(str(out))
    return out


def test_extracts_every_slide(sample_pptx, tmp_path):
    result = digest_office.extract_pptx(sample_pptx, "L1-s1", tmp_path / "assets")
    assert len(result.slides) == 2
    assert result.slides[0].title == "What is a micro:bit?"
    assert result.slides[1].title == "The LED grid"


def test_captures_speaker_notes(sample_pptx, tmp_path):
    result = digest_office.extract_pptx(sample_pptx, "L1-s1", tmp_path / "assets")
    assert "predict first" in result.slides[0].notes


def test_extracts_images_to_assets_dir(sample_pptx, tmp_path):
    assets = tmp_path / "assets"
    result = digest_office.extract_pptx(sample_pptx, "L1-s1", assets)
    assert len(result.images) == 1
    written = list(assets.glob("*.png"))
    assert len(written) == 1
    assert written[0].stat().st_size > 0


def test_image_manifest_records_source_slide(sample_pptx, tmp_path):
    result = digest_office.extract_pptx(sample_pptx, "L1-s1", tmp_path / "assets")
    assert result.images[0]["slide"] == 2


def test_warns_on_empty_slide(tmp_path):
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    src = tmp_path / "empty.pptx"
    prs.save(str(src))

    result = digest_office.extract_pptx(src, "L1-s2", tmp_path / "assets")
    assert any("empty" in w for w in result.warnings)


def test_rejects_invalid_session_id(sample_pptx, tmp_path):
    with pytest.raises(ValueError):
        digest_office.extract_pptx(sample_pptx, "L9-s9", tmp_path / "assets")


def test_image_only_slide_not_flagged_empty(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout, no text
    png = tmp_path / "dot.png"
    png.write_bytes(
        bytes.fromhex(
            "89504e470d0a1a0a0000000d4948445200000001000000010802000000907753"
            "de0000000c4944415408d76360000000020001e221bc330000000049454e44ae426082"
        )
    )
    slide.shapes.add_picture(str(png), Inches(1), Inches(1))
    src = tmp_path / "image_only.pptx"
    prs.save(str(src))

    result = digest_office.extract_pptx(src, "L1-s3", tmp_path / "assets")
    assert not any("empty" in w for w in result.warnings)


def test_warns_when_image_save_fails(sample_pptx, tmp_path, monkeypatch):
    from pptx.shapes.picture import Picture

    def broken_image(self):
        raise ValueError("corrupt blob")

    monkeypatch.setattr(Picture, "image", property(broken_image))

    result = digest_office.extract_pptx(sample_pptx, "L1-s4", tmp_path / "assets")
    assert any("failed to save image" in w for w in result.warnings)
    assert len(result.images) == 0
